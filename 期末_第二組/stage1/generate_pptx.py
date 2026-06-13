"""
Stage 1 → PPTX 簡報產生器
==========================
讀取 figures/ 與 tables/ 內容，依章節組裝成 16:9 簡報。

執行方式：
    python -m stage1.generate_pptx

輸出：
    stage1_report.pptx  （在專案根目錄）
"""

import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from stage1._common import ROOT, FIG_DIR, TBL_DIR

# ============================================================
# 顏色（與圖表配色一致）
# ============================================================
C_FRANCE = RGBColor(0x1B, 0x7A, 0x5A)
C_ROSE   = RGBColor(0xD8, 0x5A, 0x30)
C_VAR    = RGBColor(0x1B, 0x7A, 0x5A)
C_GLOBAL = RGBColor(0x2C, 0x5F, 0x8A)
C_TEXT   = RGBColor(0x3D, 0x3D, 0x3D)
C_TEXT2  = RGBColor(0x8A, 0x8A, 0x8A)
C_BG     = RGBColor(0xFA, 0xFA, 0xF7)
C_GRID   = RGBColor(0xE8, 0xE6, 0xE0)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

# 預設中文字體（macOS / Windows 都能 fallback）
ZH_FONT = 'PingFang TC'
ZH_FONT_FALLBACK = 'Microsoft JhengHei'

# 16:9 投影片
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_zh_font(run, size=None, bold=None, color=None):
    """套用中文字體 + 可選大小/粗細/顏色"""
    run.font.name = ZH_FONT
    # 中文字體 fallback
    rpr = run._r.get_or_add_rPr()
    from lxml import etree
    east_asian = etree.SubElement(rpr, '{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
    east_asian.set('typeface', ZH_FONT)
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def add_textbox(slide, left, top, width, height, text,
                size=14, bold=False, color=C_TEXT, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_zh_font(run, size=size, bold=bold, color=color)
    return tb


def add_image_centered(slide, image_path, max_w=Inches(11.5), max_h=Inches(5.6),
                        top=Inches(1.5)):
    """嵌入圖片並依長寬比縮放置中"""
    from PIL import Image
    img = Image.open(image_path)
    iw, ih = img.size
    aspect = iw / ih

    # 適配框
    target_w = max_w
    target_h = int(target_w / aspect)
    if target_h > max_h:
        target_h = max_h
        target_w = int(target_h * aspect)

    left = int((SLIDE_W - target_w) / 2)
    return slide.shapes.add_picture(image_path, left, top,
                                     width=target_w, height=target_h)


def add_color_bar(slide, color, top=Inches(0), height=Inches(0.18)):
    """投影片頂部色條"""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, top, SLIDE_W, height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def make_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


# ============================================================
# 投影片建構：標題頁
# ============================================================
def add_title_slide(prs):
    s = make_blank_slide(prs)

    # 大背景色塊（左 1/3）
    side = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(4), SLIDE_H
    )
    side.fill.solid()
    side.fill.fore_color.rgb = C_FRANCE
    side.line.fill.background()

    # 副標 "Stage 1 期中分析"
    add_textbox(s, Inches(0.6), Inches(2.5), Inches(3), Inches(0.5),
                'Stage 1 期中分析', size=14, bold=True, color=C_WHITE)

    # 主標
    add_textbox(s, Inches(4.5), Inches(2.6), Inches(8.5), Inches(1.5),
                'Provence Rosé', size=44, bold=True, color=C_FRANCE)
    add_textbox(s, Inches(4.5), Inches(3.4), Inches(8.5), Inches(1),
                'Premiumization 策略分析', size=30, bold=True, color=C_TEXT)

    # 副標
    add_textbox(s, Inches(4.5), Inches(4.4), Inches(8.5), Inches(0.6),
                '以公開貿易資料與消費者評鑑數據探索',
                size=14, color=C_TEXT2)
    add_textbox(s, Inches(4.5), Inches(4.75), Inches(8.5), Inches(0.6),
                '法國普羅旺斯 Rosé 的市場定位',
                size=14, color=C_TEXT2)

    # 資料來源
    add_textbox(s, Inches(4.5), Inches(6.3), Inches(8.5), Inches(0.4),
                '資料來源：UN Comtrade  ·  Wine Enthusiast  ·  法國海關 DGDDI',
                size=11, color=C_TEXT2)


# ============================================================
# 投影片建構：章節分隔頁
# ============================================================
def add_section_slide(prs, ch_no, title, subtitle, color):
    s = make_blank_slide(prs)
    add_color_bar(s, color)

    # 大章節編號
    add_textbox(s, Inches(0.8), Inches(1.8), Inches(2.5), Inches(2),
                f'Ch.{ch_no}', size=80, bold=True, color=color)

    # 章節標題
    add_textbox(s, Inches(3.5), Inches(2.4), Inches(9), Inches(1),
                title, size=32, bold=True, color=C_TEXT)
    # 章節副題
    add_textbox(s, Inches(3.5), Inches(3.5), Inches(9), Inches(2),
                subtitle, size=15, color=C_TEXT2)


# ============================================================
# 投影片建構：圖表頁
# ============================================================
def add_figure_slide(prs, ch_no, title, fig_path, caption=None, color=C_FRANCE):
    s = make_blank_slide(prs)
    add_color_bar(s, color)

    # 章節 + 標題
    add_textbox(s, Inches(0.5), Inches(0.35), Inches(2), Inches(0.5),
                f'Ch.{ch_no}', size=14, bold=True, color=color)
    add_textbox(s, Inches(0.5), Inches(0.65), Inches(12), Inches(0.6),
                title, size=22, bold=True, color=C_TEXT)

    # 圖
    add_image_centered(s, fig_path,
                       max_w=Inches(12), max_h=Inches(5.4),
                       top=Inches(1.4))

    # 註解
    if caption:
        add_textbox(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.5),
                    caption, size=11, color=C_TEXT2, align=PP_ALIGN.CENTER)


# ============================================================
# 投影片建構：結論頁（key findings 表格化）
# ============================================================
def add_findings_slide(prs):
    s = make_blank_slide(prs)
    add_color_bar(s, C_GLOBAL)

    add_textbox(s, Inches(0.5), Inches(0.35), Inches(2), Inches(0.5),
                'Ch.5', size=14, bold=True, color=C_GLOBAL)
    add_textbox(s, Inches(0.5), Inches(0.65), Inches(12), Inches(0.6),
                '報告核心數字一覽', size=22, bold=True, color=C_TEXT)

    df = pd.read_csv(os.path.join(TBL_DIR, 'ch5_key_findings.csv'))
    # 取每章前 4 行避免太擠
    items = []
    for ch in ['Ch.1', 'Ch.2', 'Ch.3']:
        sub = df[df['章節'] == ch].head(4)
        items.extend(sub.to_dict('records'))

    n_rows = len(items) + 1  # 標題列
    n_cols = 4

    table = s.shapes.add_table(
        n_rows, n_cols,
        Inches(0.5), Inches(1.5),
        Inches(12.3), Inches(0.4) * n_rows
    ).table

    # 欄寬
    table.columns[0].width = Inches(1.0)
    table.columns[1].width = Inches(4.0)
    table.columns[2].width = Inches(3.5)
    table.columns[3].width = Inches(3.8)

    # 標題列
    headers = ['章節', '指標', '數字', '備註']
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_FRANCE
        cell.text_frame.clear()
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = h
        set_zh_font(run, size=12, bold=True, color=C_WHITE)

    # 資料列
    for i, row in enumerate(items, start=1):
        is_alt = i % 2 == 0
        for j, key in enumerate(['章節', '指標', '數字', '備註']):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_BG if not is_alt else C_WHITE
            cell.text_frame.clear()
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j != 0 else PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(row[key])
            set_zh_font(run, size=10, bold=(j == 2),
                        color=C_TEXT if j != 2 else C_FRANCE)

    add_textbox(s, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4),
                '完整 12 張統計表詳見 tables/ 資料夾',
                size=10, color=C_TEXT2, align=PP_ALIGN.CENTER)


# ============================================================
# 投影片建構：結論文字頁
# ============================================================
def add_conclusion_slide(prs):
    s = make_blank_slide(prs)
    add_color_bar(s, C_GLOBAL)

    add_textbox(s, Inches(0.5), Inches(0.35), Inches(2), Inches(0.5),
                'Ch.5', size=14, bold=True, color=C_GLOBAL)
    add_textbox(s, Inches(0.5), Inches(0.65), Inches(12), Inches(0.6),
                '結論：三層論證、研究限制與期末展望',
                size=22, bold=True, color=C_TEXT)

    # 三大結論
    bullets = [
        ('🍷 三層數據支撐 Premiumization',
         '全球（法國 +6.1% CAGR、1.49× 溢價）→ Provence（88.06 分、$30+ 占 21.2%）→ 海關（Var 周轉 3.8 月、+4.6% 出莊 CAGR）'),
        ('📦 供需結構健康，問題在通路效率',
         'Var 的庫存周轉是 Bordeaux 的 1/6（3.8 vs 23.8 月）；瓶頸不是市場需求，而是分銷與通路。'),
        ('⚠️ Sans IG 擴張是需要關注的結構轉變',
         '6 年從 2.1% → 6.2%（年增 +0.85%, R²=0.95）。若延伸全法，恐侵蝕 Provence 的 premium 定位。'),
    ]

    y = Inches(1.6)
    for title, desc in bullets:
        add_textbox(s, Inches(0.8), y, Inches(11.7), Inches(0.5),
                    title, size=16, bold=True, color=C_FRANCE)
        add_textbox(s, Inches(0.8), y + Inches(0.45), Inches(11.7), Inches(1),
                    desc, size=12, color=C_TEXT)
        y += Inches(1.4)

    # Stage 2 預告
    add_textbox(s, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.4),
                '【期末預告】Stage 2：',
                size=14, bold=True, color=C_GLOBAL)
    add_textbox(s, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.5),
                '價格預測模型（Ridge/RF/XGBoost 比較）+ 出口市場優先排序矩陣（ML 排名）',
                size=12, color=C_TEXT)


# ============================================================
# 主流程
# ============================================================
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1. 標題頁
    add_title_slide(prs)

    # 2. Ch.1
    add_section_slide(
        prs, 1, '全球葡萄酒貿易與 Premiumization 趨勢',
        '法國「量跌、值穩、價漲」— 主要出口國 6 年市占率與單價趨勢，'
        '驗證高端化趨勢成立。',
        C_GLOBAL,
    )
    add_figure_slide(prs, 1, '主要出口國市占率（2019-2024）',
                     os.path.join(FIG_DIR, 'fig1_0_market_share.png'),
                     '法國市占率穩居第一；義大利為唯一保住量的大產區。',
                     C_GLOBAL)
    add_figure_slide(prs, 1, '2024 年出口國定位：量 × 價 × 值',
                     os.path.join(FIG_DIR, 'fig1_1_global_positioning.png'),
                     '法國在 2024 年出口均價達 1.49× 全球平均，溢價定位明確。',
                     C_GLOBAL)
    add_figure_slide(prs, 1, '法國「量跌、值穩、價漲」雙面板',
                     os.path.join(FIG_DIR, 'fig1_2_france_volume_value.png'),
                     'CAGR：量 -3.4% / 值 +2.5% / 價 +6.1%，premiumization 動力強勁。',
                     C_GLOBAL)
    add_figure_slide(prs, 1, '主要出口國單價趨勢（2019-2024）',
                     os.path.join(FIG_DIR, 'fig1_3_price_trend.png'),
                     '法國單價遠高於其他國（Welch t: p<.001***）。',
                     C_GLOBAL)

    # 3. Ch.2
    add_section_slide(
        prs, 2, 'Rosé 市場結構與 Provence 品質定位',
        'Provence 不是靠價格便宜或昂貴脫穎而出，而是靠評分與風味語彙'
        '建立差異化。',
        C_ROSE,
    )
    add_figure_slide(prs, 2, '全球 Rosé 價格帶分佈（4 大產區國）',
                     os.path.join(FIG_DIR, 'fig2_1_price_tier.png'),
                     '法國 $30+ 高端帶占 21.2%（χ² 檢定 p<.001***）。',
                     C_ROSE)
    add_figure_slide(prs, 2, '法國 11 個 Rosé 產區定位',
                     os.path.join(FIG_DIR, 'fig2_2_france_regions.png'),
                     '普羅旺斯 954 支、$22、88.06 分 — 數量第一、評分中高、價格中位。',
                     C_ROSE)
    add_figure_slide(prs, 2, 'Provence Rosé vs 非 Provence Rosé',
                     os.path.join(FIG_DIR, 'fig2_3_provence_vs_others.png'),
                     '評分顯著差異（d=+0.36, p<.001）；價格無實質差距。',
                     C_ROSE)
    add_figure_slide(prs, 2, 'Provence Rosé 風味 DNA（TF-IDF Top 15）',
                     os.path.join(FIG_DIR, 'fig2_4_provence_terms.png'),
                     '特徵詞：fruits / aftertaste / acidity / fruity / ripe / rich / crisp / texture。',
                     C_ROSE)

    # 4. Ch.3
    add_section_slide(
        prs, 3, '法國產區供需結構分析',
        'Var 的庫存周轉率僅 3.8 月，是 Bordeaux 的 1/6；Sans IG 占比 6 年從'
        '2.1% 線性升至 6.2%。',
        C_VAR,
    )
    add_figure_slide(prs, 3, '法國主要產酒區四象限定位',
                     os.path.join(FIG_DIR, 'fig3_1_quadrant.png'),
                     'Var (83) 與 BdR (13) 落在「穩定成長」象限。',
                     C_VAR)
    add_figure_slide(prs, 3, '5 大產區庫存周轉率比較',
                     os.path.join(FIG_DIR, 'fig3_2_turnover_compare.png'),
                     'Var 配對 t vs Bordeaux：t=-21.31, p<.001***。',
                     C_VAR)
    add_figure_slide(prs, 3, 'Var (83) 供需平衡',
                     os.path.join(FIG_DIR, 'fig3_3_var_supply_demand.png'),
                     '出莊量持續高於產量，2024-25 庫存從高峰下降 14% — 需求拉動。',
                     C_VAR)
    add_figure_slide(prs, 3, 'Var Sans IG 結構變遷',
                     os.path.join(FIG_DIR, 'fig3_4_var_ig_structure.png'),
                     'Sans IG 線性趨勢年增 +0.85%, R²=0.95（極強趨勢）。',
                     C_VAR)
    add_figure_slide(prs, 3, '月度季節性 + COVID 前後',
                     os.path.join(FIG_DIR, 'fig3_5_seasonal_covid.png'),
                     '出莊旺季 1-6 月；COVID 與非 COVID 期 Mann-Whitney 檢定差異不顯著。',
                     C_VAR)

    # 5. Ch.4
    add_section_slide(
        prs, 4, 'Premiumization 策略整合',
        '三層論證儀表板 + 進口市場優先排序矩陣，為 Stage 2 ML 模型奠基。',
        C_FRANCE,
    )
    add_figure_slide(prs, 4, '三層論證 KPI 儀表板',
                     os.path.join(FIG_DIR, 'fig4_1_dashboard.png'),
                     '九個 KPI 串起三層論證：全球趨勢 → Provence 品質 → 海關供需。',
                     C_FRANCE)
    add_figure_slide(prs, 4, '進口市場優先排序矩陣（2024）',
                     os.path.join(FIG_DIR, 'fig4_2_import_market.png'),
                     '高價高量 = 優先進入；高價低量 = 利基培育（Stage 2 ML 排序基礎）。',
                     C_FRANCE)

    # 6. Ch.5
    add_findings_slide(prs)
    add_conclusion_slide(prs)

    out_path = os.path.join(ROOT, 'stage1_report.pptx')
    prs.save(out_path)
    print(f'✅ PPTX 完成：{out_path}')
    print(f'   共 {len(prs.slides)} 張投影片')


if __name__ == '__main__':
    build()
